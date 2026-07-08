import argparse
import base64
import io
import json
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Dict, List, Optional

import matplotlib.pyplot as plt
import pandas as pd
import requests
from pptx import Presentation
from pptx.util import Inches


@dataclass
class AgentConfig:
    litellm_complete_url: str
    litellm_model: str
    litellm_api_key: str
    litellm_max_seconds: int


class MCPError(RuntimeError):
    pass


class TableauMCPClient:
    def __init__(self, base_url: str, timeout_seconds: int = 90) -> None:
        self.base_url = base_url.rstrip("/")
        self.timeout_seconds = timeout_seconds
        self.session_id: Optional[str] = None
        self._request_id = 1

    def _next_id(self) -> int:
        rid = self._request_id
        self._request_id += 1
        return rid

    def initialize(self) -> Dict[str, Any]:
        payload = {
            "jsonrpc": "2.0",
            "id": self._next_id(),
            "method": "initialize",
            "params": {
                "protocolVersion": "2025-03-26",
                "capabilities": {},
                "clientInfo": {"name": "tableau-dashboard-briefing-agent", "version": "1.0.0"},
            },
        }
        data, headers = self._post(payload, use_session=False)
        self.session_id = headers.get("mcp-session-id")
        if not self.session_id:
            raise MCPError("Initialize succeeded but no mcp-session-id header was returned")
        return data

    def list_tools(self) -> List[Dict[str, Any]]:
        payload = {"jsonrpc": "2.0", "id": self._next_id(), "method": "tools/list", "params": {}}
        data, _ = self._post(payload, use_session=True)
        tools = data.get("result", {}).get("tools", [])
        if not isinstance(tools, list):
            raise MCPError("Unexpected tools/list format")
        return tools

    def call_tool(self, tool_name: str, arguments: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
        payload = {
            "jsonrpc": "2.0",
            "id": self._next_id(),
            "method": "tools/call",
            "params": {"name": tool_name, "arguments": arguments or {}},
        }
        data, _ = self._post(payload, use_session=True)
        return data

    def _post(self, payload: Dict[str, Any], use_session: bool) -> (Dict[str, Any], Dict[str, str]):
        headers = {
            "Content-Type": "application/json",
            "Accept": "application/json, text/event-stream",
        }
        if use_session:
            if not self.session_id:
                raise MCPError("MCP session is missing. Initialize first.")
            headers["mcp-session-id"] = self.session_id

        response = requests.post(self.base_url, headers=headers, json=payload, timeout=self.timeout_seconds)
        if response.status_code >= 400:
            raise MCPError(f"MCP HTTP {response.status_code}: {response.text[:800]}")

        parsed = self._parse_response(response.text, expected_id=payload["id"])
        if "error" in parsed:
            raise MCPError(f"MCP RPC error: {parsed['error']}")
        return parsed, {k.lower(): v for k, v in response.headers.items()}

    @staticmethod
    def _parse_response(text: str, expected_id: int) -> Dict[str, Any]:
        stripped = text.strip()
        if not stripped:
            raise MCPError("Empty MCP response")

        if stripped.startswith("event:") or "\ndata:" in stripped:
            events: List[Dict[str, Any]] = []
            for line in stripped.splitlines():
                if line.startswith("data:"):
                    raw = line[len("data:") :].strip()
                    if not raw:
                        continue
                    try:
                        events.append(json.loads(raw))
                    except json.JSONDecodeError:
                        continue

            for event in events:
                if event.get("id") == expected_id:
                    return event

            raise MCPError("No JSON-RPC payload matching request id found in SSE response")

        return json.loads(stripped)


class LiteLLMClient:
    def __init__(self, config: AgentConfig) -> None:
        self.config = config

    def summarize(self, prompt: str) -> str:
        headers = {
            "Authorization": f"Bearer {self.config.litellm_api_key}",
            "Content-Type": "application/json",
        }
        payload = {
            "model": self.config.litellm_model,
            "messages": [
                {
                    "role": "system",
                    "content": "You are a BI analyst. Return concise insight bullets for executive slides.",
                },
                {"role": "user", "content": prompt},
            ],
            "temperature": 0.2,
        }
        response = requests.post(
            self.config.litellm_complete_url,
            headers=headers,
            json=payload,
            timeout=self.config.litellm_max_seconds,
        )
        response.raise_for_status()
        body = response.json()
        return body["choices"][0]["message"]["content"].strip()


def load_agent_config(config_path: Path) -> AgentConfig:
    data = json.loads(config_path.read_text(encoding="utf-8"))
    return AgentConfig(
        litellm_complete_url=data["LITELLM_COMPLETE_URL"],
        litellm_model=data["LITELLM_MODEL"],
        litellm_api_key=data["LITELLM_API_KEY"],
        litellm_max_seconds=int(data.get("LITELLM_MAX_SECONDS", 120)),
    )


def parse_tool_content_text(tool_response: Dict[str, Any]) -> str:
    content = tool_response.get("result", {}).get("content", [])
    if not content:
        return ""
    first = content[0]
    return str(first.get("text", ""))


def parse_tool_content_image(tool_response: Dict[str, Any]) -> Optional[bytes]:
    content = tool_response.get("result", {}).get("content", [])
    if not content:
        return None
    first = content[0]
    if first.get("type") != "image":
        return None
    raw = first.get("data")
    if not raw:
        return None
    return base64.b64decode(raw)


def list_views(client: TableauMCPClient, limit: int = 200) -> List[Dict[str, Any]]:
    resp = client.call_tool("list-views", {"limit": limit})
    text = parse_tool_content_text(resp)
    if not text:
        return []
    try:
        views = json.loads(text)
    except json.JSONDecodeError:
        return []
    if not isinstance(views, list):
        return []
    return views


def choose_view(views: List[Dict[str, Any]], dashboard_name: Optional[str], dashboard_id: Optional[str]) -> Dict[str, Any]:
    if not views:
        raise RuntimeError("No Tableau views found for this site")

    if dashboard_id:
        for v in views:
            if v.get("id") == dashboard_id:
                return v
        raise RuntimeError(f"No view found with id {dashboard_id}")

    if dashboard_name:
        dashboard_name_l = dashboard_name.lower()
        for v in views:
            if dashboard_name_l in str(v.get("name", "")).lower():
                return v
        raise RuntimeError(f"No view found matching name '{dashboard_name}'")

    print("Available dashboards/views:")
    for idx, v in enumerate(views, start=1):
        print(f"{idx}. {v.get('name')} | id={v.get('id')} | workbook={v.get('workbook', {}).get('id')}")

    selected = input("Select dashboard/view index: ").strip()
    if not selected.isdigit() or int(selected) < 1 or int(selected) > len(views):
        raise RuntimeError("Invalid selection index")
    return views[int(selected) - 1]


def parse_csv_data(csv_text: str) -> pd.DataFrame:
    if not csv_text:
        return pd.DataFrame()
    try:
        return pd.read_csv(io.StringIO(csv_text))
    except Exception:
        return pd.DataFrame()


def generate_numeric_chart(df: pd.DataFrame, output_dir: Path) -> Optional[Path]:
    if df.empty:
        return None

    numeric_cols = [c for c in df.columns if pd.api.types.is_numeric_dtype(df[c])]
    if not numeric_cols:
        # Try coercion for columns that may be percentages like '19.5%'
        candidate_cols = []
        tmp = df.copy()
        for c in df.columns:
            series = tmp[c].astype(str).str.replace("%", "", regex=False)
            coerced = pd.to_numeric(series, errors="coerce")
            if coerced.notna().sum() >= 3:
                tmp[c] = coerced
                candidate_cols.append(c)
        if not candidate_cols:
            return None
        df = tmp
        numeric_cols = candidate_cols

    value_col = numeric_cols[0]
    label_col = next((c for c in df.columns if c != value_col and not pd.api.types.is_numeric_dtype(df[c])), None)

    chart_path = output_dir / "dashboard_data_chart.png"
    plt.figure(figsize=(11, 5.5))

    if label_col:
        plot_df = df[[label_col, value_col]].dropna().head(15)
        plt.bar(plot_df[label_col].astype(str), plot_df[value_col])
        plt.xticks(rotation=45, ha="right")
        plt.xlabel(label_col)
    else:
        plot_df = df[[value_col]].dropna().head(25)
        plt.plot(range(len(plot_df)), plot_df[value_col], marker="o")
        plt.xlabel("Row")

    plt.ylabel(value_col)
    plt.title(f"Dashboard Data Snapshot: {value_col}")
    plt.tight_layout()
    plt.savefig(chart_path)
    plt.close()
    return chart_path


def build_markdown_report(
    output_dir: Path,
    selected_view: Dict[str, Any],
    summary: str,
    df: pd.DataFrame,
    mcp_tools: List[str],
) -> Path:
    md_path = output_dir / "dashboard_summary.md"

    preview = "No tabular preview available."
    if not df.empty:
        preview = df.head(12).to_markdown(index=False)

    text = [
        f"# Tableau Dashboard Briefing - {selected_view.get('name', 'Unknown')}",
        "",
        "## Dashboard Metadata",
        f"- View ID: {selected_view.get('id')}",
        f"- View Name: {selected_view.get('name')}",
        f"- Content URL: {selected_view.get('contentUrl')}",
        f"- Workbook ID: {selected_view.get('workbook', {}).get('id')}",
        "",
        "## Executive Summary",
        summary,
        "",
        "## Data Preview",
        preview,
        "",
        "## MCP Tools Available",
        ", ".join(mcp_tools),
        "",
    ]

    md_path.write_text("\n".join(text), encoding="utf-8")
    return md_path


def build_powerpoint(
    output_dir: Path,
    selected_view: Dict[str, Any],
    summary: str,
    dashboard_image_path: Optional[Path],
    chart_path: Optional[Path],
) -> Path:
    ppt = Presentation()

    # Slide 1: title
    slide = ppt.slides.add_slide(ppt.slide_layouts[0])
    slide.shapes.title.text = f"Dashboard Briefing: {selected_view.get('name', 'Dashboard')}"
    slide.placeholders[1].text = "Generated by Tableau Dashboard Briefing Agent"

    # Slide 2: concise insight bullets
    slide = ppt.slides.add_slide(ppt.slide_layouts[1])
    slide.shapes.title.text = "Executive Insights"
    tf = slide.placeholders[1].text_frame
    tf.clear()
    lines = [ln.strip("- ").strip() for ln in summary.splitlines() if ln.strip()]
    compact = lines[:6] if lines else ["No summary available from LLM."]
    for i, line in enumerate(compact):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line[:180]

    # Slide 3: dashboard image
    slide = ppt.slides.add_slide(ppt.slide_layouts[5])
    slide.shapes.title.text = "Dashboard Visual"
    if dashboard_image_path and dashboard_image_path.exists():
        slide.shapes.add_picture(str(dashboard_image_path), Inches(0.5), Inches(1.2), width=Inches(12.2))

    # Slide 4: analytical chart
    slide = ppt.slides.add_slide(ppt.slide_layouts[5])
    slide.shapes.title.text = "Analytical Chart"
    if chart_path and chart_path.exists():
        slide.shapes.add_picture(str(chart_path), Inches(0.8), Inches(1.3), width=Inches(11.8))

    ppt_path = output_dir / "dashboard_briefing.pptx"
    ppt.save(str(ppt_path))
    return ppt_path


def run_agent(
    config_path: Path,
    mcp_url: str,
    dashboard_name: Optional[str],
    dashboard_id: Optional[str],
) -> None:
    output_dir = Path("outputs") / "dashboard_briefing"
    output_dir.mkdir(parents=True, exist_ok=True)

    config = load_agent_config(config_path)
    llm = LiteLLMClient(config)
    client = TableauMCPClient(mcp_url)

    client.initialize()
    tools = client.list_tools()
    tool_names = sorted([t.get("name", "") for t in tools])

    required = {"list-views", "get-view", "get-view-data", "get-view-image"}
    missing = [name for name in required if name not in set(tool_names)]
    if missing:
        raise RuntimeError(f"Missing required Tableau MCP tools: {missing}")

    views = list_views(client, limit=300)
    selected_view = choose_view(views, dashboard_name=dashboard_name, dashboard_id=dashboard_id)
    view_id = selected_view.get("id")

    view_resp = client.call_tool("get-view", {"viewId": view_id})
    view_data_resp = client.call_tool("get-view-data", {"viewId": view_id})
    view_image_resp = client.call_tool("get-view-image", {"viewId": view_id, "format": "PNG", "width": 1400, "height": 800})

    view_json = parse_tool_content_text(view_resp)
    csv_text = parse_tool_content_text(view_data_resp)
    image_bytes = parse_tool_content_image(view_image_resp)

    raw_dir = output_dir / "raw"
    raw_dir.mkdir(parents=True, exist_ok=True)
    (raw_dir / "selected_view.json").write_text(view_json if view_json else json.dumps(selected_view, indent=2), encoding="utf-8")
    (raw_dir / "selected_view_data.csv").write_text(csv_text, encoding="utf-8")

    dashboard_image_path: Optional[Path] = None
    if image_bytes:
        dashboard_image_path = output_dir / "dashboard_image.png"
        dashboard_image_path.write_bytes(image_bytes)

    df = parse_csv_data(csv_text)
    chart_path = generate_numeric_chart(df, output_dir)

    sample_data = df.head(20).to_markdown(index=False) if not df.empty else "No tabular data available"
    prompt = (
        "Create concise executive insights for one slide.\n"
        "Requirements: 5-7 bullets, analytical style, short lines, no long paragraphs.\n"
        f"Dashboard name: {selected_view.get('name')}\n"
        f"View metadata: {json.dumps(selected_view)}\n"
        f"Data sample:\n{sample_data}\n"
    )

    try:
        summary = llm.summarize(prompt)
    except Exception as exc:
        summary = f"- LLM summary unavailable: {exc}\n- Use the data preview for manual insights."

    md_path = build_markdown_report(output_dir, selected_view, summary, df, tool_names)
    ppt_path = build_powerpoint(output_dir, selected_view, summary, dashboard_image_path, chart_path)

    artifacts = {
        "selected_view": selected_view,
        "markdown_report": str(md_path),
        "powerpoint": str(ppt_path),
        "dashboard_image": str(dashboard_image_path) if dashboard_image_path else None,
        "chart_image": str(chart_path) if chart_path else None,
        "raw_data_csv": str(raw_dir / "selected_view_data.csv"),
    }
    (output_dir / "artifacts.json").write_text(json.dumps(artifacts, indent=2), encoding="utf-8")

    print("Dashboard briefing generation complete.")
    print(f"Selected view: {selected_view.get('name')} ({selected_view.get('id')})")
    print(f"Markdown report: {md_path}")
    print(f"PowerPoint: {ppt_path}")
    if dashboard_image_path:
        print(f"Dashboard image: {dashboard_image_path}")
    if chart_path:
        print(f"Analytical chart: {chart_path}")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Read Tableau dashboard and generate concise MD + PPT analytical briefing")
    parser.add_argument("--config", default="config/config.json", help="Path to Litellm config JSON")
    parser.add_argument("--mcp-url", default="http://localhost:7777/tableau-mcp", help="Tableau MCP URL")
    parser.add_argument("--dashboard-name", default=None, help="View/dashboard name (substring match)")
    parser.add_argument("--dashboard-id", default=None, help="Exact view ID")

    args = parser.parse_args()
    run_agent(
        config_path=Path(args.config),
        mcp_url=args.mcp_url,
        dashboard_name=args.dashboard_name,
        dashboard_id=args.dashboard_id,
    )
