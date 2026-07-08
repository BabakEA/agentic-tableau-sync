import datetime as dt
import base64
import io
import json
import os
import re
from pathlib import Path
from typing import Any, Dict, List, Literal, Optional, Tuple

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
import requests
import yfinance as yf
from fastapi import FastAPI, HTTPException
from pydantic import BaseModel, Field

try:
    import tableauserverclient as TSC
except Exception:
    TSC = None

try:
    import pantab
except Exception:
    pantab = None


class MCPError(RuntimeError):
    pass


class PublishError(RuntimeError):
    pass


class TableauMCPClient:
    def __init__(self, base_url: str, timeout_seconds: int = 90) -> None:
        self.base_url = base_url.rstrip("/")
        self.timeout_seconds = timeout_seconds
        self.session_id: Optional[str] = None
        self._request_id = 1

    def _next_id(self) -> int:
        rid = self._request_id
        self._request_id += 1
        return rid

    def initialize(self) -> Dict[str, Any]:
        payload = {
            "jsonrpc": "2.0",
            "id": self._next_id(),
            "method": "initialize",
            "params": {
                "protocolVersion": "2025-03-26",
                "capabilities": {},
                "clientInfo": {"name": "stock-dashboard-api", "version": "1.0.0"},
            },
        }
        data, headers = self._post(payload, use_session=False)
        self.session_id = headers.get("mcp-session-id")
        if not self.session_id:
            raise MCPError("MCP initialize did not return mcp-session-id")
        return data

    def list_tools(self) -> List[Dict[str, Any]]:
        payload = {"jsonrpc": "2.0", "id": self._next_id(), "method": "tools/list", "params": {}}
        data, _ = self._post(payload, use_session=True)
        tools = data.get("result", {}).get("tools", [])
        if not isinstance(tools, list):
            raise MCPError("Unexpected tools/list response format")
        return tools

    def call_tool(self, tool_name: str, arguments: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
        payload = {
            "jsonrpc": "2.0",
            "id": self._next_id(),
            "method": "tools/call",
            "params": {"name": tool_name, "arguments": arguments or {}},
        }
        data, _ = self._post(payload, use_session=True)
        return data

    def _post(self, payload: Dict[str, Any], use_session: bool) -> Tuple[Dict[str, Any], Dict[str, str]]:
        headers = {
            "Content-Type": "application/json",
            "Accept": "application/json, text/event-stream",
        }
        if use_session:
            if not self.session_id:
                raise MCPError("Session ID not set. Call initialize first.")
            headers["mcp-session-id"] = self.session_id

        response = requests.post(self.base_url, headers=headers, json=payload, timeout=self.timeout_seconds)
        if response.status_code >= 400:
            raise MCPError(f"MCP HTTP {response.status_code}: {response.text[:600]}")

        parsed = self._parse_response(response.text, expected_id=payload["id"])
        if "error" in parsed:
            raise MCPError(f"MCP RPC error: {parsed['error']}")
        return parsed, {k.lower(): v for k, v in response.headers.items()}

    @staticmethod
    def _parse_response(text: str, expected_id: int) -> Dict[str, Any]:
        stripped = text.strip()
        if not stripped:
            raise MCPError("Empty MCP response body")

        if stripped.startswith("event:") or "\ndata:" in stripped:
            events: List[Dict[str, Any]] = []
            for line in stripped.splitlines():
                if not line.startswith("data:"):
                    continue
                raw = line[len("data:") :].strip()
                if not raw:
                    continue
                try:
                    events.append(json.loads(raw))
                except json.JSONDecodeError:
                    continue

            for event in events:
                if event.get("id") == expected_id:
                    return event

            raise MCPError("No JSON-RPC response matching request id found in SSE stream")

        return json.loads(stripped)


def parse_text_content(tool_response: Dict[str, Any]) -> str:
    content = tool_response.get("result", {}).get("content", [])
    if not content:
        return ""
    return str(content[0].get("text", ""))


def load_tableau_publish_config(env_path: Path = Path("env.list")) -> Dict[str, str]:
    values: Dict[str, str] = {}
    if env_path.exists():
        for raw_line in env_path.read_text(encoding="utf-8").splitlines():
            line = raw_line.strip()
            if not line or line.startswith("#") or "=" not in line:
                continue
            key, val = line.split("=", 1)
            values[key.strip()] = val.strip().strip('"').strip("'")

    server = os.getenv("TABLEAU_SERVER", values.get("SERVER", ""))
    site = os.getenv("TABLEAU_SITE_NAME", values.get("SITE_NAME", ""))
    pat_name = os.getenv("TABLEAU_PAT_NAME", values.get("PAT_NAME", ""))
    pat_value = os.getenv("TABLEAU_PAT_VALUE", values.get("PAT_VALUE", ""))

    return {
        "server": server,
        "site_name": site,
        "pat_name": pat_name,
        "pat_value": pat_value,
    }


def load_litellm_config(config_path: Path = Path("config/config.json")) -> Dict[str, Any]:
    data = json.loads(config_path.read_text(encoding="utf-8"))
    return {
        "complete_url": data["LITELLM_COMPLETE_URL"],
        "model": data["LITELLM_MODEL"],
        "api_key": data["LITELLM_API_KEY"],
        "max_seconds": int(data.get("LITELLM_MAX_SECONDS", 120)),
    }


class LiteLLMClient:
    def __init__(self, config: Dict[str, Any]) -> None:
        self.config = config

    def chat(self, messages: List[Dict[str, Any]], temperature: float = 0.2) -> str:
        headers = {
            "Authorization": f"Bearer {self.config['api_key']}",
            "Content-Type": "application/json",
        }
        payload = {
            "model": self.config["model"],
            "messages": messages,
            "temperature": temperature,
        }
        response = requests.post(
            self.config["complete_url"],
            headers=headers,
            json=payload,
            timeout=self.config["max_seconds"],
        )
        response.raise_for_status()
        body = response.json()
        content = body["choices"][0]["message"]["content"]
        if isinstance(content, str):
            return content.strip()
        if isinstance(content, list):
            parts: List[str] = []
            for item in content:
                if isinstance(item, dict) and item.get("type") == "text":
                    parts.append(str(item.get("text", "")))
            return "\n".join(part for part in parts if part).strip()
        return str(content).strip()


def canonicalize_ticker(raw_ticker: str) -> Optional[str]:
    token = raw_ticker.strip().upper().replace(" ", "")
    if not token:
        return None
    aliases = {
        "SP500": "^GSPC",
        "S&P500": "^GSPC",
        "S&P_500": "^GSPC",
        "S&P": "^GSPC",
        "GSPC": "^GSPC",
        "SPX": "^GSPC",
        "^SPX": "^GSPC",
        "^GSPC": "^GSPC",
    }
    return aliases.get(token, token)


def parse_requested_tickers(raw_ticker: str, include_sp500: bool, sp500_sample_size: int) -> Tuple[List[str], str]:
    tickers: List[str] = []
    for token in raw_ticker.replace(";", ",").split(","):
        canonical = canonicalize_ticker(token)
        if canonical and canonical not in tickers:
            tickers.append(canonical)

    primary = next((ticker for ticker in tickers if ticker != "^GSPC"), None)
    if not primary:
        raise RuntimeError("Please provide at least one stock ticker, for example AAPL")

    if "^GSPC" not in tickers:
        tickers.append("^GSPC")

    if include_sp500:
        for symbol in fetch_sp500_tickers(limit=sp500_sample_size):
            if symbol not in tickers:
                tickers.append(symbol)

    return tickers, primary


def build_dataframe_profile(df: pd.DataFrame) -> str:
    if df.empty:
        return "No tabular data was extracted."

    lines = [f"Rows: {len(df)}", f"Columns: {', '.join(df.columns.astype(str).tolist()[:20])}"]

    numeric_cols = [col for col in df.columns if pd.api.types.is_numeric_dtype(df[col])]
    if numeric_cols:
        lines.append("Numeric summary:")
        for col in numeric_cols[:6]:
            series = df[col].dropna()
            if series.empty:
                continue
            lines.append(
                f"- {col}: mean={series.mean():.4f}, median={series.median():.4f}, min={series.min():.4f}, max={series.max():.4f}"
            )

    categorical_cols = [col for col in df.columns if col not in numeric_cols]
    if categorical_cols:
        lines.append("Categorical summary:")
        for col in categorical_cols[:4]:
            top_values = df[col].astype(str).value_counts().head(3).to_dict()
            lines.append(f"- {col}: top values {top_values}")

    time_col = next((col for col in df.columns if any(key in str(col).lower() for key in ["date", "time", "month", "year"])), None)
    if time_col:
        time_series = pd.to_datetime(df[time_col], errors="coerce").dropna()
        if not time_series.empty:
            lines.append(f"Time range: {time_series.min()} to {time_series.max()}")

    try:
        sample = df.head(8).to_markdown(index=False)
    except Exception:
        sample = df.head(8).to_csv(index=False)
    lines.append("Sample rows:")
    lines.append(sample)
    return "\n".join(lines)


def split_lines_for_slide(text: str, limit: int = 8) -> List[str]:
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    cleaned = [line.lstrip("- ").strip() for line in lines]
    return cleaned[:limit]


def run_agentic_view_analysis(
    llm: LiteLLMClient,
    view_name: str,
    workbook_name: str,
    image_path: Optional[Path],
    df: pd.DataFrame,
    include_kid_friendly: bool,
) -> Dict[str, str]:
    profile = build_dataframe_profile(df)

    visual_prompt = (
        f"You are a BI analytics agent analyzing a Tableau view image.\n"
        f"Workbook: {workbook_name}\n"
        f"View: {view_name}\n"
        "Look at the chart carefully and describe what is visually present, the main trend, anomalies, comparisons, and what a presenter should mention.\n"
        "Return concise bullets only.\n"
        f"Data profile for grounding:\n{profile}"
    )

    visual_observation: str
    try:
        if image_path and image_path.exists():
            image_b64 = base64.b64encode(image_path.read_bytes()).decode("utf-8")
            visual_observation = llm.chat(
                [
                    {
                        "role": "system",
                        "content": "You are an expert analytics agent. Be precise, concise, and evidence-driven.",
                    },
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": visual_prompt},
                            {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{image_b64}"}},
                        ],
                    },
                ]
            )
        else:
            raise RuntimeError("No image available for vision analysis")
    except Exception:
        visual_observation = llm.chat(
            [
                {
                    "role": "system",
                    "content": "You are an expert analytics agent. Be precise, concise, and evidence-driven.",
                },
                {
                    "role": "user",
                    "content": (
                        visual_prompt
                        + "\nThe image could not be analyzed directly. Use the data profile to infer likely insights, and say when a point is data-derived."
                    ),
                },
            ]
        )

    synthesis_prompt = (
        f"Using the observations below and the dataframe profile, generate presentation-ready content for workbook '{workbook_name}' view '{view_name}'.\n"
        "Return three sections exactly with bullets:\n"
        "Executive Bullets:\n- ...\n"
        "Data Notes:\n- ...\n"
        "Kid Explanation:\n- ...\n"
        f"Include kid explanation: {'yes' if include_kid_friendly else 'no'}\n\n"
        f"Visual observations:\n{visual_observation}\n\n"
        f"Data profile:\n{profile}"
    )

    synthesis = llm.chat(
        [
            {
                "role": "system",
                "content": "You are a senior BI storytelling agent creating concise, presentation-quality content.",
            },
            {"role": "user", "content": synthesis_prompt},
        ]
    )

    return {
        "profile": profile,
        "visual_observation": visual_observation,
        "synthesis": synthesis,
    }


def run_workbook_summary_agent(llm: LiteLLMClient, workbook_name: str, view_analyses: List[Dict[str, Any]]) -> str:
    compact = []
    for item in view_analyses:
        compact.append(f"View: {item['view_name']}\n{item.get('analysis_text', '')}")
    prompt = (
        f"Create an executive summary for Tableau workbook '{workbook_name}'.\n"
        "Return 5-7 concise bullets for the title/overview slide. Focus on the most decision-relevant insights across all views.\n\n"
        + "\n\n".join(compact[:12])
    )
    return llm.chat(
        [
            {
                "role": "system",
                "content": "You are a senior BI storytelling agent creating concise executive summaries.",
            },
            {"role": "user", "content": prompt},
        ]
    )


def _ensure_tsc_available() -> None:
    if TSC is None:
        raise PublishError("tableauserverclient is not installed. Install with: pip install tableauserverclient")


def _ensure_pantab_available() -> None:
    if pantab is None:
        raise PublishError("pantab is not installed. Install with: pip install pantab")


def ensure_project(server: Any, project_name: str) -> Any:
    all_projects, _ = server.projects.get()
    for project in all_projects:
        if project.name == project_name:
            return project

    item = TSC.ProjectItem(name=project_name)
    return server.projects.create(item)


def publish_stock_datasource(
    run_dir: Path,
    dashboard_name: str,
    project_name: str,
    publish_cfg: Dict[str, str],
    stock_df: pd.DataFrame,
) -> Dict[str, Any]:
    _ensure_tsc_available()
    _ensure_pantab_available()

    missing = [k for k, v in publish_cfg.items() if not v]
    if missing:
        raise PublishError(f"Missing Tableau publish config keys: {missing}")

    hyper_path = run_dir / f"{dashboard_name}_stock_data.hyper"
    # Single-table extract. Tableau can consume it as a published datasource.
    pantab.frame_to_hyper(stock_df, str(hyper_path), table="Extract")

    auth = TSC.PersonalAccessTokenAuth(
        token_name=publish_cfg["pat_name"],
        personal_access_token=publish_cfg["pat_value"],
        site_id=publish_cfg["site_name"],
    )
    server = TSC.Server(publish_cfg["server"], use_server_version=True)

    with server.auth.sign_in(auth):
        project = ensure_project(server, project_name)
        datasource_name = f"{dashboard_name}_data"
        ds_item = TSC.DatasourceItem(project_id=project.id, name=datasource_name)
        published = server.datasources.publish(ds_item, str(hyper_path), mode=TSC.Server.PublishMode.Overwrite)
        return {
            "datasource_id": published.id,
            "datasource_name": published.name,
            "project_id": project.id,
            "project_name": project.name,
            "hyper_path": str(hyper_path),
        }


def publish_workbook_template(
    dashboard_name: str,
    project_name: str,
    publish_cfg: Dict[str, str],
    workbook_template_path: str,
) -> Dict[str, Any]:
    _ensure_tsc_available()

    template_path = Path(workbook_template_path)
    if not template_path.exists():
        raise PublishError(f"Workbook template not found: {workbook_template_path}")

    auth = TSC.PersonalAccessTokenAuth(
        token_name=publish_cfg["pat_name"],
        personal_access_token=publish_cfg["pat_value"],
        site_id=publish_cfg["site_name"],
    )
    server = TSC.Server(publish_cfg["server"], use_server_version=True)

    with server.auth.sign_in(auth):
        project = ensure_project(server, project_name)
        wb_item = TSC.WorkbookItem(project_id=project.id, name=dashboard_name, show_tabs=True)
        published = server.workbooks.publish(wb_item, str(template_path), mode=TSC.Server.PublishMode.Overwrite)
        return {
            "workbook_id": published.id,
            "workbook_name": published.name,
            "project_id": project.id,
            "project_name": project.name,
            "template_path": str(template_path),
        }


def clone_and_publish_workbook_from_server(
    dashboard_name: str,
    project_name: str,
    publish_cfg: Dict[str, str],
    source_workbook_id: str,
    run_dir: Path,
) -> Dict[str, Any]:
    _ensure_tsc_available()

    auth = TSC.PersonalAccessTokenAuth(
        token_name=publish_cfg["pat_name"],
        personal_access_token=publish_cfg["pat_value"],
        site_id=publish_cfg["site_name"],
    )
    server = TSC.Server(publish_cfg["server"], use_server_version=True)

    with server.auth.sign_in(auth):
        # Download existing workbook from Tableau, then republish under a new name.
        downloaded_path = Path(server.workbooks.download(source_workbook_id, filepath=str(run_dir), include_extract=True))
        project = ensure_project(server, project_name)
        wb_item = TSC.WorkbookItem(project_id=project.id, name=dashboard_name, show_tabs=True)
        published = server.workbooks.publish(wb_item, str(downloaded_path), mode=TSC.Server.PublishMode.Overwrite)
        return {
            "workbook_id": published.id,
            "workbook_name": published.name,
            "project_id": project.id,
            "project_name": project.name,
            "source_workbook_id": source_workbook_id,
            "downloaded_template": str(downloaded_path),
        }


def normalize_csv_text(text: str) -> str:
    # Some MCP responses return CSV as a JSON-escaped quoted string.
    stripped = text.strip()
    if stripped.startswith('"') and stripped.endswith('"'):
        try:
            return json.loads(stripped)
        except json.JSONDecodeError:
            return text
    return text


def fetch_sp500_tickers(limit: int = 30) -> List[str]:
    try:
        response = requests.get(
            "https://en.wikipedia.org/wiki/List_of_S%26P_500_companies",
            headers={"User-Agent": "Mozilla/5.0"},
            timeout=30,
        )
        response.raise_for_status()
        table = pd.read_html(response.text)[0]
        tickers = table["Symbol"].astype(str).tolist()
        tickers = [t.replace(".", "-") for t in tickers]
        return tickers[:limit]
    except Exception:
        return []


def fetch_price_history(ticker: str, period: str = "6mo", interval: str = "1d") -> pd.DataFrame:
    history = yf.Ticker(ticker).history(period=period, interval=interval, auto_adjust=False)
    if history.empty:
        raise ValueError(f"No price history returned for {ticker}")

    history = history.reset_index().rename(columns={"Date": "timestamp"})
    if "timestamp" not in history.columns:
        first_col = history.columns[0]
        history = history.rename(columns={first_col: "timestamp"})
    history["ticker"] = ticker
    return history


def infer_timestamp_column(df: pd.DataFrame) -> str:
    for candidate in ["timestamp", "Datetime", "Date", "index"]:
        if candidate in df.columns:
            return candidate
    return df.columns[0]


def fetch_options_snapshot(ticker: str) -> pd.DataFrame:
    yf_ticker = yf.Ticker(ticker)
    expirations = yf_ticker.options
    if not expirations:
        return pd.DataFrame(
            [{"ticker": ticker, "expiration": None, "avg_call_iv": np.nan, "avg_put_iv": np.nan, "calls": 0, "puts": 0}]
        )

    nearest = expirations[0]
    chain = yf_ticker.option_chain(nearest)
    calls = chain.calls if chain.calls is not None else pd.DataFrame()
    puts = chain.puts if chain.puts is not None else pd.DataFrame()

    call_iv = float(calls["impliedVolatility"].dropna().mean()) if not calls.empty else np.nan
    put_iv = float(puts["impliedVolatility"].dropna().mean()) if not puts.empty else np.nan

    return pd.DataFrame(
        [
            {
                "ticker": ticker,
                "expiration": nearest,
                "avg_call_iv": call_iv,
                "avg_put_iv": put_iv,
                "calls": int(len(calls)),
                "puts": int(len(puts)),
            }
        ]
    )


def compute_metrics(prices: pd.DataFrame) -> pd.DataFrame:
    rows: List[Dict[str, Any]] = []
    for ticker, frame in prices.groupby("ticker"):
        frame = frame.sort_values("timestamp").copy()
        frame["daily_return"] = frame["Close"].pct_change()
        mean_ret = float(frame["daily_return"].mean(skipna=True))
        vol = float(frame["daily_return"].std(skipna=True))
        close_last = float(frame["Close"].iloc[-1])
        close_first = float(frame["Close"].iloc[0])
        perf = (close_last / close_first - 1.0) if close_first else np.nan
        rows.append(
            {
                "ticker": ticker,
                "period_return": perf,
                "avg_daily_return": mean_ret,
                "daily_volatility": vol,
                "last_close": close_last,
            }
        )
    return pd.DataFrame(rows).sort_values("period_return", ascending=False)


def build_charts(prices: pd.DataFrame, metrics: pd.DataFrame, output_dir: Path) -> List[Path]:
    output_dir.mkdir(parents=True, exist_ok=True)
    chart_paths: List[Path] = []

    plt.figure(figsize=(12, 6))
    for ticker, frame in prices.groupby("ticker"):
        frame = frame.sort_values("timestamp")
        norm = frame["Close"] / frame["Close"].iloc[0]
        plt.plot(frame["timestamp"], norm, label=ticker)
    plt.title("Normalized Price Performance")
    plt.xlabel("Date")
    plt.ylabel("Normalized Price")
    plt.legend()
    plt.tight_layout()
    p1 = output_dir / "normalized_performance.png"
    plt.savefig(p1)
    plt.close()
    chart_paths.append(p1)

    plt.figure(figsize=(10, 6))
    plt.scatter(metrics["daily_volatility"], metrics["period_return"])
    for _, row in metrics.iterrows():
        plt.annotate(row["ticker"], (row["daily_volatility"], row["period_return"]))
    plt.title("Return vs Volatility")
    plt.xlabel("Daily Volatility")
    plt.ylabel("Period Return")
    plt.tight_layout()
    p2 = output_dir / "return_vs_volatility.png"
    plt.savefig(p2)
    plt.close()
    chart_paths.append(p2)

    return chart_paths


def build_moving_average_chart(
    ticker_prices: pd.DataFrame,
    spx_prices: pd.DataFrame,
    output_dir: Path,
    ma_window: int,
) -> Path:
    ticker_prices = ticker_prices.sort_values("timestamp").copy()
    spx_prices = spx_prices.sort_values("timestamp").copy()

    ticker_prices["ma"] = ticker_prices["Close"].rolling(ma_window, min_periods=1).mean()
    spx_prices["base_norm"] = spx_prices["Close"] / spx_prices["Close"].iloc[0]
    ticker_prices["asset_norm"] = ticker_prices["Close"] / ticker_prices["Close"].iloc[0]

    out = output_dir / "moving_average_vs_sp500.png"
    plt.figure(figsize=(12, 6))
    plt.plot(ticker_prices["timestamp"], ticker_prices["Close"], label="Asset Close", alpha=0.7)
    plt.plot(ticker_prices["timestamp"], ticker_prices["ma"], label=f"MA({ma_window})", linewidth=2)
    plt.plot(spx_prices["timestamp"], spx_prices["base_norm"] * ticker_prices["Close"].iloc[0], label="S&P 500 Baseline", linestyle="--")
    plt.title("Price, Moving Average, and S&P 500 Baseline")
    plt.xlabel("Time")
    plt.ylabel("Price")
    plt.legend()
    plt.tight_layout()
    plt.savefig(out)
    plt.close()
    return out


def build_hourly_signal_chart(prices: pd.DataFrame, output_dir: Path) -> Tuple[Path, Path]:
    df = prices.sort_values("timestamp").copy()
    df["timestamp"] = pd.to_datetime(df["timestamp"], errors="coerce")
    df = df.dropna(subset=["timestamp"]) 
    if df.empty:
        raise ValueError("No valid timestamps for hourly signal chart")

    df["hour"] = df["timestamp"].dt.hour
    df["month"] = df["timestamp"].dt.to_period("M").astype(str)
    df["ret"] = df["Close"].pct_change()

    # Approximate buy/sell signal from average return by hour.
    hourly = (
        df.groupby("hour", as_index=False)
        .agg(avg_return=("ret", "mean"), avg_volume=("Volume", "mean"))
        .fillna(0.0)
    )
    hourly["signal"] = np.where(hourly["avg_return"] >= 0, "buy", "sell")

    # Requested "average hour per month in one week" style rollup using first 7 days.
    week_start = df["timestamp"].min().normalize()
    week_end = week_start + pd.Timedelta(days=7)
    week = df[(df["timestamp"] >= week_start) & (df["timestamp"] < week_end)].copy()
    week_month_hour = (
        week.groupby(["month", "hour"], as_index=False)
        .agg(avg_return=("ret", "mean"), avg_volume=("Volume", "mean"))
        .fillna(0.0)
    )

    p1 = output_dir / "hourly_buy_sell_signal.png"
    plt.figure(figsize=(12, 6))
    colors = ["green" if s == "buy" else "red" for s in hourly["signal"]]
    plt.bar(hourly["hour"], hourly["avg_return"], color=colors)
    plt.axhline(0, color="black", linewidth=1)
    plt.title("Average Return by Hour (Buy/Sell Signal)")
    plt.xlabel("Hour of Day")
    plt.ylabel("Average Return")
    plt.tight_layout()
    plt.savefig(p1)
    plt.close()

    p2 = output_dir / "hourly_volume_week_snapshot.png"
    plt.figure(figsize=(12, 6))
    if not week_month_hour.empty:
        for month, grp in week_month_hour.groupby("month"):
            plt.plot(grp["hour"], grp["avg_volume"], marker="o", label=month)
        plt.legend()
    plt.title("Average Volume by Hour (First Week Snapshot)")
    plt.xlabel("Hour of Day")
    plt.ylabel("Average Volume")
    plt.tight_layout()
    plt.savefig(p2)
    plt.close()

    hourly.to_csv(output_dir / "hourly_signal_table.csv", index=False)
    week_month_hour.to_csv(output_dir / "week_month_hour_table.csv", index=False)
    return p1, p2


def extract_views(client: TableauMCPClient, limit: int = 200) -> List[Dict[str, Any]]:
    response = client.call_tool("list-views", {"limit": limit})
    text = parse_text_content(response)
    if not text:
        return []
    try:
        views = json.loads(text)
    except json.JSONDecodeError:
        return []
    if not isinstance(views, list):
        return []
    return views


def extract_views_filtered(client: TableauMCPClient, limit: int = 200, filter_expr: Optional[str] = None) -> List[Dict[str, Any]]:
    args: Dict[str, Any] = {"limit": limit}
    if filter_expr:
        args["filter"] = filter_expr
    response = client.call_tool("list-views", args)
    text = parse_text_content(response)
    if not text:
        return []
    try:
        views = json.loads(text)
    except json.JSONDecodeError:
        return []
    if not isinstance(views, list):
        return []
    return views


def extract_datasources(client: TableauMCPClient, limit: int = 200, filter_expr: Optional[str] = None) -> List[Dict[str, Any]]:
    args: Dict[str, Any] = {"limit": limit}
    if filter_expr:
        args["filter"] = filter_expr
    response = client.call_tool("list-datasources", args)
    text = parse_text_content(response)
    if not text:
        return []
    try:
        items = json.loads(text)
    except json.JSONDecodeError:
        return []
    if not isinstance(items, list):
        return []
    return items


def dashboard_exists(client: TableauMCPClient, dashboard_name: str) -> bool:
    tools = {t.get("name") for t in client.list_tools() if isinstance(t, dict)}

    if "search-content" in tools:
        try:
            response = client.call_tool(
                "search-content",
                {"terms": dashboard_name, "filter": {"contentTypes": ["view", "workbook"]}, "limit": 30},
            )
            if dashboard_name.lower() in json.dumps(response).lower():
                return True
        except Exception:
            pass

    if "list-views" in tools:
        try:
            for view in extract_views(client, limit=300):
                if dashboard_name.lower() in str(view.get("name", "")).lower():
                    return True
        except Exception:
            pass

    return False


class StockDashboardRequest(BaseModel):
    dashboard_name: str = Field(default="Sstock_market")
    ticker: str = Field(default="AAPL")
    period: Literal["5d", "7d", "1mo", "3mo", "6mo", "1y", "2y", "5y", "10y"] = "6mo"
    interval: Literal["1m", "2m", "5m", "15m", "30m", "60m", "90m", "1h", "1d", "1wk", "1mo"] = "1d"
    moving_average_window: int = Field(default=20, ge=2, le=300)
    include_sp500: bool = False
    sp500_sample_size: int = 10
    mcp_url: str = "http://localhost:7777/tableau-mcp"
    publish: bool = True
    publish_project_name: str = "AI_Agent_Dashboards"
    workbook_template_path: Optional[str] = None
    source_workbook_id: Optional[str] = None


class DashboardBriefingRequest(BaseModel):
    dashboard_name: str
    mcp_url: str = "http://localhost:7777/tableau-mcp"
    include_kid_friendly: bool = True
    workbook_name: Optional[str] = None
    workbook_id: Optional[str] = None


class ViewBriefingRequest(BaseModel):
    view_id: Optional[str] = None
    view_name: Optional[str] = None
    workbook_id: Optional[str] = None
    mcp_url: str = "http://localhost:7777/tableau-mcp"
    include_kid_friendly: bool = True


class ChatTurn(BaseModel):
    role: Literal["user", "assistant"]
    content: str


class ChatRequest(BaseModel):
    message: str
    history: List[ChatTurn] = Field(default_factory=list)
    mcp_url: str = "http://localhost:7777/tableau-mcp"


def build_workbook_link(server_url: str, site_name: str, workbook_id: str) -> str:
    server = server_url.rstrip("/")
    return f"{server}/#/site/{site_name}/workbooks/{workbook_id}"


def build_view_link(server_url: str, site_name: str, content_url: Optional[str]) -> Optional[str]:
    if not content_url:
        return None
    server = server_url.rstrip("/")
    # content_url usually looks like: sample/sheets/ExecutiveSummary
    parts = content_url.split("/")
    if len(parts) >= 3 and parts[1] == "sheets":
        workbook_slug = parts[0]
        view_slug = parts[2]
        return f"{server}/#/site/{site_name}/views/{workbook_slug}/{view_slug}"
    return f"{server}/#/site/{site_name}/views/{content_url}"


def summarize_view_data(df: pd.DataFrame, include_kid_friendly: bool) -> List[str]:
    if df.empty:
        lines = ["No tabular data extracted for this view."]
        if include_kid_friendly:
            lines.append("Kid version: This page is like a picture board without a score table behind it.")
        return lines

    numeric_cols = [c for c in df.columns if pd.api.types.is_numeric_dtype(df[c])]
    lines: List[str] = [f"Rows analyzed: {len(df)}"]
    if numeric_cols:
        top_cols = numeric_cols[:3]
        for c in top_cols:
            series = df[c].dropna()
            if not series.empty:
                lines.append(f"{c}: avg={series.mean():.3f}, min={series.min():.3f}, max={series.max():.3f}")
    else:
        lines.append("No numeric measures detected in this view dataset.")

    if include_kid_friendly:
        lines.append("Kid version: Bigger bars/lines usually mean more, smaller ones mean less.")
    return lines


def resolve_target_view(client: TableauMCPClient, req: ViewBriefingRequest) -> Dict[str, Any]:
    views = extract_views(client, limit=500)
    candidates = views

    if req.workbook_id:
        candidates = [v for v in candidates if ((v.get("workbook") or {}).get("id") == req.workbook_id)]

    if req.view_id:
        for v in candidates:
            if v.get("id") == req.view_id:
                return v
        raise RuntimeError(f"View not found for id: {req.view_id}")

    if req.view_name:
        for v in candidates:
            if req.view_name.lower() in str(v.get("name", "")).lower():
                return v
        raise RuntimeError(f"View not found for name: {req.view_name}")

    raise RuntimeError("Please provide view_id or view_name")


DEFAULT_SOURCE_WORKBOOK_ID = "8b646cec-4f8c-4e24-900a-c052ffb83332"


def find_uuid(text: str) -> Optional[str]:
    match = re.search(r"\b[0-9a-fA-F]{8}-[0-9a-fA-F]{4}-[0-9a-fA-F]{4}-[0-9a-fA-F]{4}-[0-9a-fA-F]{12}\b", text)
    return match.group(0) if match else None


def heuristic_chat_plan(message: str) -> Optional[Dict[str, Any]]:
    text = message.strip()
    lowered = text.lower()
    report_words = ["powerpoint", "ppt", "report", "summaris", "summariz"]

    if "show me" in lowered and "dashboard" in lowered and "view" in lowered:
        return {"action": "list_catalog"}
    if "show" in lowered and "dashboard" in lowered:
        return {"action": "list_dashboards"}
    if "show" in lowered and "view" in lowered:
        return {"action": "list_views"}
    if "datasource" in lowered:
        return {"action": "list_datasources"}
    if any(word in lowered for word in report_words) and "workbook" in lowered:
        maybe_id = find_uuid(text)
        if maybe_id:
            return {"action": "report_dashboard", "workbook_id": maybe_id}
        match = re.search(r"workbook\s+([A-Za-z0-9_\- ]+?)(?:\s+and\s+|$)", text, re.IGNORECASE)
        workbook_name = match.group(1).strip() if match else None
        return {"action": "report_dashboard", "workbook_name": workbook_name or text}
    if any(word in lowered for word in report_words) and "view" in lowered:
        maybe_id = find_uuid(text)
        if maybe_id:
            return {"action": "report_view", "view_id": maybe_id}
        match = re.search(r"view\s+([A-Za-z0-9_\- ]+?)(?:\s+and\s+|$)", text, re.IGNORECASE)
        view_name = match.group(1).strip() if match else None
        return {"action": "report_view", "view_name": view_name or text}
    return None


def llm_chat_plan(llm: LiteLLMClient, message: str, history: List[ChatTurn]) -> Dict[str, Any]:
    history_text = "\n".join([f"{turn.role}: {turn.content}" for turn in history[-8:]])
    prompt = (
        "You are an orchestration agent for a Tableau analytics backend. Return JSON only.\n"
        "Schema:\n"
        "{\n"
        '  "action": "list_dashboards|list_views|list_datasources|list_catalog|report_dashboard|report_view|create_stock_dashboard|help",\n'
        '  "dashboard_name": null,\n'
        '  "workbook_name": null,\n'
        '  "workbook_id": null,\n'
        '  "view_name": null,\n'
        '  "view_id": null,\n'
        '  "ticker": null,\n'
        '  "period": "3mo",\n'
        '  "interval": "1h",\n'
        '  "moving_average_window": 20,\n'
        '  "publish": true,\n'
        '  "include_kid_friendly": true\n'
        "}\n"
        "Rules:\n"
        "- If user asks for dashboards and views, use list_catalog.\n"
        "- If user asks to create a stock dashboard and no workbook template is given, use source workbook clone fallback handled by backend.\n"
        f"Conversation history:\n{history_text}\n\n"
        f"User message:\n{message}"
    )
    raw = llm.chat(
        [
            {"role": "system", "content": "You are a strict JSON planner. Return valid JSON only."},
            {"role": "user", "content": prompt},
        ],
        temperature=0,
    )
    start = raw.find("{")
    end = raw.rfind("}")
    if start >= 0 and end > start:
        try:
            return json.loads(raw[start : end + 1])
        except Exception:
            pass
    return {"action": "help"}


def collect_artifact_paths(obj: Any) -> List[str]:
    paths: List[str] = []

    def walk(value: Any) -> None:
        if isinstance(value, dict):
            for item in value.values():
                walk(item)
        elif isinstance(value, list):
            for item in value:
                walk(item)
        elif isinstance(value, str):
            lower = value.lower()
            if any(lower.endswith(ext) for ext in [".pptx", ".md", ".csv", ".png", ".json", ".txt", ".hyper", ".twb", ".twbx"]):
                paths.append(value)

    walk(obj)
    unique: List[str] = []
    seen = set()
    for path in paths:
        if path not in seen:
            seen.add(path)
            unique.append(path)
    return unique


def summarize_chat_result(action: str, result: Dict[str, Any]) -> str:
    if action == "list_dashboards":
        names = [item.get("name") for item in result.get("dashboards", [])[:10]]
        return f"Found {result.get('count', 0)} dashboards. Top results: {', '.join([n for n in names if n])}."
    if action == "list_views":
        names = [item.get("name") for item in result.get("views", [])[:10]]
        return f"Found {result.get('count', 0)} views. Top results: {', '.join([n for n in names if n])}."
    if action == "list_datasources":
        names = [item.get("name") for item in result.get("datasources", [])[:10]]
        return f"Found {result.get('count', 0)} datasources. Top results: {', '.join([n for n in names if n])}."
    if action == "list_catalog":
        return (
            f"Found {result.get('dashboards', {}).get('count', 0)} dashboards and "
            f"{result.get('views', {}).get('count', 0)} views."
        )
    if action == "create_stock_dashboard":
        dashboard = result.get("dashboard", {})
        return f"Stock dashboard job finished. Status: {dashboard.get('status')}. Publish status: {dashboard.get('publish_status')}."
    if action in {"report_dashboard", "report_view"}:
        return "Report and PowerPoint were generated. Use the files below to download them."
    return "I processed your request."


app = FastAPI(title="Tableau Stock Agents API", version="1.0.0")


@app.get("/health")
def health() -> Dict[str, str]:
    return {"status": "ok"}


@app.get("/dashboards")
def list_dashboards(
    mcp_url: str = "http://localhost:7777/tableau-mcp",
    limit: int = 100,
    dashboard_only: bool = True,
) -> Dict[str, Any]:
    try:
        client = TableauMCPClient(mcp_url)
        client.initialize()
        filter_expr = "sheetType:eq:dashboard" if dashboard_only else None
        views = extract_views_filtered(client, limit=limit, filter_expr=filter_expr)
        items = [
            {
                "id": v.get("id"),
                "name": v.get("name"),
                "contentUrl": v.get("contentUrl"),
                "workbookId": (v.get("workbook") or {}).get("id"),
            }
            for v in views
        ]
        return {"count": len(items), "dashboard_only": dashboard_only, "dashboards": items}
    except Exception as exc:
        raise HTTPException(status_code=500, detail=str(exc)) from exc


@app.get("/views")
def list_views_endpoint(
    mcp_url: str = "http://localhost:7777/tableau-mcp",
    limit: int = 200,
    dashboard_only: bool = False,
    workbook_id: Optional[str] = None,
) -> Dict[str, Any]:
    try:
        client = TableauMCPClient(mcp_url)
        client.initialize()

        filter_expr = "sheetType:eq:dashboard" if dashboard_only else None
        views = extract_views_filtered(client, limit=limit, filter_expr=filter_expr)
        if workbook_id:
            views = [v for v in views if ((v.get("workbook") or {}).get("id") == workbook_id)]

        items = [
            {
                "id": v.get("id"),
                "name": v.get("name"),
                "contentUrl": v.get("contentUrl"),
                "workbookId": (v.get("workbook") or {}).get("id"),
                "projectId": (v.get("project") or {}).get("id"),
            }
            for v in views
        ]
        return {
            "count": len(items),
            "dashboard_only": dashboard_only,
            "workbook_id": workbook_id,
            "views": items,
        }
    except Exception as exc:
        raise HTTPException(status_code=500, detail=str(exc)) from exc


@app.get("/datasources")
def list_datasources_endpoint(
    mcp_url: str = "http://localhost:7777/tableau-mcp",
    limit: int = 200,
    name: Optional[str] = None,
) -> Dict[str, Any]:
    try:
        client = TableauMCPClient(mcp_url)
        client.initialize()

        filter_expr = f"name:eq:{name}" if name else None
        sources = extract_datasources(client, limit=limit, filter_expr=filter_expr)
        items = [
            {
                "id": item.get("id"),
                "name": item.get("name"),
                "contentUrl": item.get("contentUrl"),
                "projectName": item.get("projectName"),
                "ownerName": item.get("ownerName"),
            }
            for item in sources
        ]
        return {"count": len(items), "datasources": items}
    except Exception as exc:
        raise HTTPException(status_code=500, detail=str(exc)) from exc


@app.get("/meta/options")
def get_dropdown_options() -> Dict[str, Any]:
    return {
        "periods": ["5d", "7d", "1mo", "3mo", "6mo", "1y", "2y", "5y", "10y"],
        "intervals": ["1m", "2m", "5m", "15m", "30m", "60m", "90m", "1h", "1d", "1wk", "1mo"],
        "defaults": {
            "period": "6mo",
            "interval": "1d",
            "moving_average_window": 20,
            "ticker": "AAPL",
            "dashboard_name": "Sstock_market",
        },
    }


@app.post("/agent/chat")
def agent_chat(request: ChatRequest) -> Dict[str, Any]:
    try:
        llm = LiteLLMClient(load_litellm_config())
        plan = heuristic_chat_plan(request.message) or llm_chat_plan(llm, request.message, request.history)
        action = plan.get("action", "help")

        if action == "list_dashboards":
            result = list_dashboards(mcp_url=request.mcp_url, limit=50, dashboard_only=True)
        elif action == "list_views":
            result = list_views_endpoint(mcp_url=request.mcp_url, limit=100, dashboard_only=False, workbook_id=plan.get("workbook_id"))
        elif action == "list_datasources":
            result = list_datasources_endpoint(mcp_url=request.mcp_url, limit=100, name=plan.get("dashboard_name"))
        elif action == "list_catalog":
            result = {
                "dashboards": list_dashboards(mcp_url=request.mcp_url, limit=50, dashboard_only=True),
                "views": list_views_endpoint(mcp_url=request.mcp_url, limit=100, dashboard_only=False, workbook_id=plan.get("workbook_id")),
            }
        elif action == "report_dashboard":
            result = generate_dashboard_powerpoint(
                DashboardBriefingRequest(
                    dashboard_name=plan.get("dashboard_name") or request.message,
                    mcp_url=request.mcp_url,
                    include_kid_friendly=bool(plan.get("include_kid_friendly", True)),
                    workbook_name=plan.get("workbook_name"),
                    workbook_id=plan.get("workbook_id"),
                )
            )
        elif action == "report_view":
            result = generate_view_briefing_comprehensive(
                ViewBriefingRequest(
                    view_id=plan.get("view_id"),
                    view_name=plan.get("view_name") if not plan.get("view_id") else None,
                    workbook_id=plan.get("workbook_id"),
                    mcp_url=request.mcp_url,
                    include_kid_friendly=bool(plan.get("include_kid_friendly", True)),
                )
            )
        elif action == "create_stock_dashboard":
            result = generate_stock_market_dashboard(
                StockDashboardRequest(
                    dashboard_name=plan.get("dashboard_name") or "Sstock_market",
                    ticker=plan.get("ticker") or request.message,
                    period=plan.get("period") or "3mo",
                    interval=plan.get("interval") or "1h",
                    moving_average_window=int(plan.get("moving_average_window") or 20),
                    include_sp500=False,
                    sp500_sample_size=10,
                    mcp_url=request.mcp_url,
                    publish=bool(plan.get("publish", True)),
                    publish_project_name="AI_Agent_Dashboards",
                    source_workbook_id=plan.get("source_workbook_id") or DEFAULT_SOURCE_WORKBOOK_ID,
                )
            )
        else:
            result = {
                "help": [
                    "Show me your dashboards",
                    "Show me your views",
                    "Summarize workbook sample and generate PPT",
                    "Generate stock dashboard for AAPL over 3 months hourly",
                    "Create report for view 3883cb98-6a7d-4738-9da7-5522566a26f4",
                ]
            }

        artifacts = collect_artifact_paths(result)
        return {
            "action": action,
            "plan": plan,
            "message": summarize_chat_result(action, result),
            "result": result,
            "artifacts": artifacts,
        }
    except Exception as exc:
        raise HTTPException(status_code=500, detail=str(exc)) from exc


@app.post("/dashboards/stock-market")
def generate_stock_market_dashboard(request: StockDashboardRequest) -> Dict[str, Any]:
    try:
        output_root = Path("outputs") / "stock_market_api"
        run_dir = output_root / dt.datetime.now().strftime("run_%Y%m%d_%H%M%S")
        run_dir.mkdir(parents=True, exist_ok=True)

        client = TableauMCPClient(request.mcp_url)
        init_data = client.initialize()
        tools = client.list_tools()
        tool_names = sorted([t.get("name", "") for t in tools if isinstance(t, dict)])

        tickers, primary_ticker = parse_requested_tickers(
            request.ticker,
            include_sp500=request.include_sp500,
            sp500_sample_size=request.sp500_sample_size,
        )

        price_frames: List[pd.DataFrame] = []
        option_frames: List[pd.DataFrame] = []
        skipped: List[Dict[str, str]] = []
        for ticker in tickers:
            try:
                if ticker == "^GSPC":
                    try:
                        price_frames.append(fetch_price_history(ticker, period=request.period, interval=request.interval))
                    except Exception:
                        spy_df = fetch_price_history("SPY", period=request.period, interval=request.interval)
                        spy_df["ticker"] = "SPY"
                        price_frames.append(spy_df)
                else:
                    price_frames.append(fetch_price_history(ticker, period=request.period, interval=request.interval))
                option_frames.append(fetch_options_snapshot(ticker))
            except Exception as exc:
                skipped.append({"ticker": ticker, "reason": str(exc)})

        if not price_frames:
            raise RuntimeError("No market data returned for provided tickers")

        prices = pd.concat(price_frames, ignore_index=True)
        prices["timestamp"] = pd.to_datetime(prices["timestamp"], errors="coerce")
        prices = prices.dropna(subset=["timestamp"]) 
        options = pd.concat(option_frames, ignore_index=True) if option_frames else pd.DataFrame()
        metrics = compute_metrics(prices)
        charts = build_charts(prices, metrics, run_dir)

        # Required analytics: moving average + hourly buy/sell recommendations.
        primary = prices[prices["ticker"] == primary_ticker].copy()
        baseline = prices[prices["ticker"].isin(["^GSPC", "SPY"])].copy()
        if primary.empty or baseline.empty:
            raise RuntimeError("Primary ticker or S&P 500 baseline data is missing")

        ma_chart = build_moving_average_chart(primary, baseline, run_dir, ma_window=request.moving_average_window)
        hourly_chart, hourly_volume_chart = build_hourly_signal_chart(primary, run_dir)
        charts.extend([ma_chart, hourly_chart, hourly_volume_chart])

        prices_path = run_dir / "stock_prices.csv"
        metrics_path = run_dir / "stock_metrics.csv"
        options_path = run_dir / "options_snapshot.csv"
        prices.to_csv(prices_path, index=False)
        metrics.to_csv(metrics_path, index=False)
        if not options.empty:
            options.to_csv(options_path, index=False)

        exists = dashboard_exists(client, request.dashboard_name)

        creation_supported = any("create" in (name or "").lower() for name in tool_names)
        publish_status = "not_requested"
        publish_note = ""
        publish_results: Dict[str, Any] = {}
        if request.publish:
            try:
                publish_cfg = load_tableau_publish_config()

                # Publish the generated stock dataset as a Tableau datasource via REST/TSC.
                publish_results["datasource"] = publish_stock_datasource(
                    run_dir=run_dir,
                    dashboard_name=request.dashboard_name,
                    project_name=request.publish_project_name,
                    publish_cfg=publish_cfg,
                    stock_df=prices,
                )

                # Optional: publish a workbook template so a visible dashboard object appears.
                template_error: Optional[str] = None
                workbook_published = False

                if request.workbook_template_path:
                    template_path = Path(request.workbook_template_path)
                    if template_path.exists():
                        publish_results["workbook"] = publish_workbook_template(
                            dashboard_name=request.dashboard_name,
                            project_name=request.publish_project_name,
                            publish_cfg=publish_cfg,
                            workbook_template_path=request.workbook_template_path,
                        )
                        publish_status = "datasource_and_workbook_published"
                        publish_note = "Datasource and workbook template published successfully."
                        workbook_published = True
                    else:
                        template_error = f"Workbook template not found: {request.workbook_template_path}"

                if not workbook_published and request.source_workbook_id:
                    try:
                        publish_results["workbook"] = clone_and_publish_workbook_from_server(
                            dashboard_name=request.dashboard_name,
                            project_name=request.publish_project_name,
                            publish_cfg=publish_cfg,
                            source_workbook_id=request.source_workbook_id,
                            run_dir=run_dir,
                        )
                        publish_status = "datasource_and_cloned_workbook_published"
                        publish_note = "Datasource published and workbook cloned from existing Tableau workbook ID."
                        if template_error:
                            publish_note += f" Template fallback used because: {template_error}."
                        workbook_published = True
                    except Exception as clone_exc:
                        if template_error:
                            raise RuntimeError(f"{template_error}. Source workbook clone also failed: {clone_exc}") from clone_exc
                        raise

                if not workbook_published:
                    publish_status = "datasource_published_no_workbook"
                    publish_note = template_error or (
                        "Datasource published successfully. To create a visible dashboard/workbook object, provide workbook_template_path or source_workbook_id."
                    )
            except Exception as exc:
                publish_status = "publish_failed"
                publish_note = str(exc)

        if publish_results.get("workbook"):
            status = "workbook_published"
        elif publish_results.get("datasource"):
            status = "datasource_published"
        else:
            status = "exists" if exists else ("create_supported" if creation_supported else "spec_generated")

        spec = {
            "dashboardName": request.dashboard_name,
            "generatedAt": dt.datetime.now(dt.timezone.utc).isoformat(),
            "status": status,
            "action": "publish_or_create_dashboard",
            "recommendedSheets": [
                "Normalized Price Performance",
                "Return vs Volatility",
                "Option IV Snapshot",
            ],
        }
        spec_path = run_dir / "tableau_dashboard_spec.json"
        spec_path.write_text(json.dumps(spec, indent=2), encoding="utf-8")

        report = {
            "initialized": init_data,
            "mcp_session_id": client.session_id,
            "mcp_tools": tool_names,
            "dashboard": {
                "name": request.dashboard_name,
                "exists": exists,
                "creation_supported": creation_supported,
                "status": status,
                "publish_requested": request.publish,
                "publish_status": publish_status,
                "publish_note": publish_note,
                "publish_results": publish_results,
                "note": "Data package generated. Publish/create depends on available Tableau write tools.",
            },
            "skipped_tickers": skipped,
            "inputs": {
                "ticker": request.ticker,
                "resolved_tickers": tickers,
                "primary_ticker": primary_ticker,
                "period": request.period,
                "interval": request.interval,
                "moving_average_window": request.moving_average_window,
            },
            "artifacts": {
                "run_dir": str(run_dir),
                "prices_csv": str(prices_path),
                "metrics_csv": str(metrics_path),
                "options_csv": str(options_path) if options_path.exists() else None,
                "charts": [str(p) for p in charts],
                "dashboard_spec": str(spec_path),
            },
        }

        publish_cfg_for_links = load_tableau_publish_config()
        server_url = publish_cfg_for_links.get("server", "")
        site_name = publish_cfg_for_links.get("site_name", "")

        if publish_results.get("workbook"):
            wb_id = publish_results["workbook"].get("workbook_id")
            if wb_id and server_url and site_name:
                report["dashboard"]["workbook_link"] = build_workbook_link(server_url, site_name, wb_id)

                # Try to add first few view links.
                try:
                    wb_resp = client.call_tool("get-workbook", {"workbookId": wb_id})
                    wb_text = parse_text_content(wb_resp)
                    wb_data = json.loads(wb_text).get("data", {}) if wb_text else {}
                    views = (wb_data.get("views") or {}).get("view", [])
                    report["dashboard"]["view_links"] = [
                        build_view_link(server_url, site_name, v.get("contentUrl")) for v in views[:10]
                    ]
                except Exception:
                    report["dashboard"]["view_links"] = []

        report_path = run_dir / "agent_report.json"
        report_path.write_text(json.dumps(report, indent=2), encoding="utf-8")
        report["report_path"] = str(report_path)
        return report
    except Exception as exc:
        raise HTTPException(status_code=500, detail=str(exc)) from exc


@app.post("/dashboards/briefing")
def generate_dashboard_powerpoint(request: DashboardBriefingRequest) -> Dict[str, Any]:
    try:
        run_dir = Path("outputs") / "dashboard_briefing_api" / dt.datetime.now().strftime("run_%Y%m%d_%H%M%S")
        run_dir.mkdir(parents=True, exist_ok=True)

        publish_cfg = load_tableau_publish_config()
        server_url = publish_cfg.get("server", "https://10ax.online.tableau.com")
        site_name = publish_cfg.get("site_name", "")
        llm = LiteLLMClient(load_litellm_config())

        client = TableauMCPClient(request.mcp_url)
        client.initialize()
        views = extract_views(client, limit=500)

        # Resolve workbook target.
        target_workbook_id: Optional[str] = request.workbook_id
        target_workbook_name: Optional[str] = request.workbook_name

        if not target_workbook_id:
            if request.workbook_name:
                wb_resp = client.call_tool("list-workbooks", {"limit": 200, "filter": f"name:eq:{request.workbook_name}"})
                wb_text = parse_text_content(wb_resp)
                wb_items = json.loads(wb_text) if wb_text else []
                if wb_items:
                    target_workbook_id = wb_items[0].get("id")
                    target_workbook_name = wb_items[0].get("name")
            else:
                # Fallback: infer workbook from a matching view name.
                selected = None
                for view in views:
                    if request.dashboard_name.lower() in str(view.get("name", "")).lower():
                        selected = view
                        break
                if selected:
                    target_workbook_id = (selected.get("workbook") or {}).get("id")
                    target_workbook_name = request.dashboard_name

        if not target_workbook_id:
            raise RuntimeError("Could not resolve workbook. Pass workbook_id or workbook_name.")

        wb_resp = client.call_tool("get-workbook", {"workbookId": target_workbook_id})
        wb_text = parse_text_content(wb_resp)
        wb_data = json.loads(wb_text).get("data", {}) if wb_text else {}
        workbook_name = wb_data.get("name") or target_workbook_name or request.dashboard_name

        wb_views = (wb_data.get("views") or {}).get("view", [])
        if not wb_views:
            # fallback to list-views filtered by workbookId if needed
            wb_views = [v for v in views if ((v.get("workbook") or {}).get("id") == target_workbook_id)]

        if not wb_views:
            raise RuntimeError("No views found for workbook")

        raw_dir = run_dir / "raw"
        images_dir = run_dir / "images"
        raw_dir.mkdir(parents=True, exist_ok=True)
        images_dir.mkdir(parents=True, exist_ok=True)

        all_view_summaries: List[Dict[str, Any]] = []
        markdown_lines: List[str] = [f"# Workbook Briefing - {workbook_name}", ""]

        for idx, view in enumerate(wb_views, start=1):
            view_id = view.get("id")
            view_name = view.get("name", f"View {idx}")
            view_content_url = view.get("contentUrl")
            data_resp = client.call_tool("get-view-data", {"viewId": view_id})
            image_resp = client.call_tool(
                "get-view-image",
                {"viewId": view_id, "format": "PNG", "width": 1400, "height": 800},
            )

            data_text = normalize_csv_text(parse_text_content(data_resp))
            df = pd.read_csv(io.StringIO(data_text)) if data_text else pd.DataFrame()
            csv_path = raw_dir / f"{idx:02d}_{view_name.replace(' ', '_')}.csv"
            if not df.empty:
                df.to_csv(csv_path, index=False)

            image_content = image_resp.get("result", {}).get("content", [])
            image_path = images_dir / f"{idx:02d}_{view_name.replace(' ', '_')}.png"
            if image_content and image_content[0].get("type") == "image" and image_content[0].get("data"):
                image_path.write_bytes(base64.b64decode(image_content[0]["data"]))

            agent_result = run_agentic_view_analysis(
                llm=llm,
                view_name=view_name,
                workbook_name=workbook_name,
                image_path=image_path if image_path.exists() else None,
                df=df,
                include_kid_friendly=request.include_kid_friendly,
            )
            summary_lines = split_lines_for_slide(agent_result["synthesis"], limit=8)
            view_link = build_view_link(server_url, site_name, view_content_url)
            analysis_path = raw_dir / f"{idx:02d}_{view_name.replace(' ', '_')}_analysis.txt"
            analysis_path.write_text(
                "\n\n".join(
                    [
                        "Visual Observation:\n" + agent_result["visual_observation"],
                        "Synthesis:\n" + agent_result["synthesis"],
                        "Profile:\n" + agent_result["profile"],
                    ]
                ),
                encoding="utf-8",
            )

            all_view_summaries.append(
                {
                    "view_id": view_id,
                    "view_name": view_name,
                    "view_link": view_link,
                    "csv": str(csv_path) if csv_path.exists() else None,
                    "image": str(image_path) if image_path.exists() else None,
                    "analysis_file": str(analysis_path),
                    "analysis_text": agent_result["synthesis"],
                    "summary": summary_lines,
                }
            )

            markdown_lines.append(f"## {view_name}")
            if view_link:
                markdown_lines.append(f"- Link: {view_link}")
            markdown_lines.extend([f"- {line}" for line in summary_lines])
            markdown_lines.append("")

        workbook_link = build_workbook_link(server_url, site_name, target_workbook_id)
        workbook_summary = run_workbook_summary_agent(llm, workbook_name, all_view_summaries)
        markdown_lines.insert(2, f"- Workbook link: {workbook_link}")
        markdown_lines.insert(3, f"- Views covered: {len(all_view_summaries)}")
        markdown_lines.insert(5, "## Executive Summary")
        markdown_lines.insert(6, workbook_summary)

        md_path = run_dir / "dashboard_briefing.md"
        md_path.write_text("\n".join(markdown_lines), encoding="utf-8")

        # Generate PowerPoint with concise text.
        from pptx import Presentation
        from pptx.util import Inches

        ppt = Presentation()
        s1 = ppt.slides.add_slide(ppt.slide_layouts[0])
        s1.shapes.title.text = f"Workbook Briefing: {workbook_name}"
        s1.placeholders[1].text = f"Views analyzed: {len(all_view_summaries)}"

        s2 = ppt.slides.add_slide(ppt.slide_layouts[1])
        s2.shapes.title.text = "Executive Summary"
        tf = s2.placeholders[1].text_frame
        tf.clear()
        headline = split_lines_for_slide(workbook_summary, limit=7)
        for i, line in enumerate(headline):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line[:180]

        for item in all_view_summaries:
            slide = ppt.slides.add_slide(ppt.slide_layouts[5])
            slide.shapes.title.text = item["view_name"]
            img = item.get("image")
            if img and Path(img).exists():
                slide.shapes.add_picture(str(img), Inches(0.4), Inches(1.2), width=Inches(8.6))

            # Right-side concise bullet summary.
            textbox = slide.shapes.add_textbox(Inches(9.2), Inches(1.2), Inches(3.8), Inches(5.5))
            tframe = textbox.text_frame
            tframe.clear()
            lines = item.get("summary", [])[:6]
            if item.get("view_link"):
                lines.insert(0, f"Link: {item['view_link']}")
            for i, line in enumerate(lines):
                p = tframe.paragraphs[0] if i == 0 else tframe.add_paragraph()
                p.text = str(line)[:170]

        ppt_path = run_dir / "dashboard_briefing.pptx"
        ppt.save(str(ppt_path))

        return {
            "workbook": {
                "name": workbook_name,
                "id": target_workbook_id,
                "link": workbook_link,
                "views_count": len(all_view_summaries),
            },
            "artifacts": {
                "run_dir": str(run_dir),
                "markdown": str(md_path),
                "powerpoint": str(ppt_path),
                "images_dir": str(images_dir),
                "raw_dir": str(raw_dir),
            },
            "views": all_view_summaries,
        }
    except Exception as exc:
        raise HTTPException(status_code=500, detail=str(exc)) from exc


@app.post("/views/briefing-comprehensive")
def generate_view_briefing_comprehensive(request: ViewBriefingRequest) -> Dict[str, Any]:
    try:
        run_dir = Path("outputs") / "view_briefing_api" / dt.datetime.now().strftime("run_%Y%m%d_%H%M%S")
        run_dir.mkdir(parents=True, exist_ok=True)
        charts_dir = run_dir / "charts"
        charts_dir.mkdir(parents=True, exist_ok=True)

        publish_cfg = load_tableau_publish_config()
        server_url = publish_cfg.get("server", "https://10ax.online.tableau.com")
        site_name = publish_cfg.get("site_name", "")
        llm = LiteLLMClient(load_litellm_config())

        client = TableauMCPClient(request.mcp_url)
        client.initialize()

        selected = resolve_target_view(client, request)
        view_id = selected.get("id")
        view_name = selected.get("name")
        workbook_id = (selected.get("workbook") or {}).get("id")

        data_resp = client.call_tool("get-view-data", {"viewId": view_id})
        image_resp = client.call_tool("get-view-image", {"viewId": view_id, "format": "PNG", "width": 1600, "height": 900})

        data_text = normalize_csv_text(parse_text_content(data_resp))
        df = pd.read_csv(io.StringIO(data_text)) if data_text else pd.DataFrame()
        data_csv_path = run_dir / "view_data.csv"
        if not df.empty:
            df.to_csv(data_csv_path, index=False)

        image_content = image_resp.get("result", {}).get("content", [])
        image_path = run_dir / "view_image.png"
        if image_content and image_content[0].get("type") == "image" and image_content[0].get("data"):
            image_path.write_bytes(base64.b64decode(image_content[0]["data"]))

        numeric_cols = [c for c in df.columns if pd.api.types.is_numeric_dtype(df[c])]
        ts_col = next((c for c in df.columns if any(k in c.lower() for k in ["date", "time", "month", "year"])), None)

        chart_paths: List[Path] = []
        if numeric_cols:
            # Top means bar chart
            top_col = numeric_cols[0]
            top_path = charts_dir / "summary_stats_bar.png"
            summary = df[numeric_cols].describe().T[["mean", "min", "max"]].head(8)
            plt.figure(figsize=(11, 6))
            summary["mean"].plot(kind="bar", color="#2a9d8f")
            plt.title("Mean Value by Numeric Field")
            plt.ylabel("Mean")
            plt.tight_layout()
            plt.savefig(top_path)
            plt.close()
            chart_paths.append(top_path)

            # Distribution chart
            hist_path = charts_dir / "distribution_hist.png"
            plt.figure(figsize=(11, 6))
            df[top_col].dropna().head(2000).hist(bins=30, color="#457b9d")
            plt.title(f"Distribution of {top_col}")
            plt.xlabel(top_col)
            plt.ylabel("Frequency")
            plt.tight_layout()
            plt.savefig(hist_path)
            plt.close()
            chart_paths.append(hist_path)

            # Trend chart if time-like column available.
            if ts_col:
                trend_path = charts_dir / "timeseries_trend.png"
                tdf = df[[ts_col, top_col]].copy()
                tdf[ts_col] = pd.to_datetime(tdf[ts_col], errors="coerce")
                tdf = tdf.dropna().sort_values(ts_col).head(5000)
                if not tdf.empty:
                    plt.figure(figsize=(11, 6))
                    plt.plot(tdf[ts_col], tdf[top_col], color="#e76f51")
                    plt.title(f"Trend of {top_col} over {ts_col}")
                    plt.xlabel(ts_col)
                    plt.ylabel(top_col)
                    plt.tight_layout()
                    plt.savefig(trend_path)
                    plt.close()
                    chart_paths.append(trend_path)

        view_link = build_view_link(server_url, site_name, selected.get("contentUrl"))
        workbook_link = build_workbook_link(server_url, site_name, workbook_id) if workbook_id else None

        agent_result = run_agentic_view_analysis(
            llm=llm,
            view_name=view_name,
            workbook_name=workbook_id or "unknown-workbook",
            image_path=image_path if image_path.exists() else None,
            df=df,
            include_kid_friendly=request.include_kid_friendly,
        )
        summary_lines = split_lines_for_slide(agent_result["synthesis"], limit=10)
        analysis_path = run_dir / "view_analysis.txt"
        analysis_path.write_text(
            "\n\n".join(
                [
                    "Visual Observation:\n" + agent_result["visual_observation"],
                    "Synthesis:\n" + agent_result["synthesis"],
                    "Profile:\n" + agent_result["profile"],
                ]
            ),
            encoding="utf-8",
        )

        md_lines = [
            f"# Comprehensive View Briefing - {view_name}",
            "",
            f"- View ID: {view_id}",
            f"- Workbook ID: {workbook_id}",
            f"- View Link: {view_link}",
            f"- Workbook Link: {workbook_link}",
            "",
            "## Key Points",
            *[f"- {line}" for line in summary_lines],
            "",
            "## Agent Narrative",
            agent_result["synthesis"],
        ]
        md_path = run_dir / "view_briefing.md"
        md_path.write_text("\n".join(md_lines), encoding="utf-8")

        from pptx import Presentation
        from pptx.util import Inches

        ppt = Presentation()
        s1 = ppt.slides.add_slide(ppt.slide_layouts[0])
        s1.shapes.title.text = f"Comprehensive Briefing: {view_name}"
        s1.placeholders[1].text = "Automated analytics from Tableau view data"

        s2 = ppt.slides.add_slide(ppt.slide_layouts[1])
        s2.shapes.title.text = "Executive Summary"
        tf = s2.placeholders[1].text_frame
        tf.clear()
        top_lines = summary_lines[:8]
        for i, line in enumerate(top_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line[:180]

        s3 = ppt.slides.add_slide(ppt.slide_layouts[5])
        s3.shapes.title.text = "Original Tableau View"
        if image_path.exists():
            s3.shapes.add_picture(str(image_path), Inches(0.5), Inches(1.2), width=Inches(12))

        for idx, cp in enumerate(chart_paths, start=1):
            slide = ppt.slides.add_slide(ppt.slide_layouts[5])
            slide.shapes.title.text = f"Analytical Chart {idx}"
            slide.shapes.add_picture(str(cp), Inches(0.6), Inches(1.2), width=Inches(11.8))

        ppt_path = run_dir / "view_briefing_comprehensive.pptx"
        ppt.save(str(ppt_path))

        return {
            "view": {
                "id": view_id,
                "name": view_name,
                "link": view_link,
                "workbook_id": workbook_id,
                "workbook_link": workbook_link,
            },
            "artifacts": {
                "run_dir": str(run_dir),
                "markdown": str(md_path),
                "powerpoint": str(ppt_path),
                "view_image": str(image_path) if image_path.exists() else None,
                "data_csv": str(data_csv_path) if data_csv_path.exists() else None,
                "analysis_text": str(analysis_path),
                "charts": [str(p) for p in chart_paths],
            },
            "summary": summary_lines,
        }
    except Exception as exc:
        raise HTTPException(status_code=500, detail=str(exc)) from exc


if __name__ == "__main__":
    import uvicorn

    uvicorn.run("agent_api:app", host="0.0.0.0", port=8010, reload=False)
